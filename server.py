import os
import io
import base64
from typing import Optional, List

from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

# Document parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from PIL import Image

# Optional OCR - will fallback gracefully if tesseract not installed
try:
    import pytesseract  # type: ignore
    TESS_AVAILABLE = True
except Exception:
    TESS_AVAILABLE = False

# Google Gemini
import google.generativeai as genai


SUPPORTED_DOC_TYPES = {".pdf", ".docx", ".pptx", ".txt"}
SUPPORTED_IMG_TYPES = {".png", ".jpg", ".jpeg"}
SUPPORTED_AUDIO_TYPES = {"audio/webm", "audio/wav", "audio/mpeg", "audio/mp4", "audio/ogg"}


def getenv_strict(name: str) -> str:
    value = os.getenv(name)
    if not value:
        raise RuntimeError(f"Missing required environment variable: {name}")
    return value


GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "AIzaSyC_TIp6ibTOksHM2mXCzBbYWXhbyjXGe6Y")
if GOOGLE_API_KEY:
    genai.configure(api_key=GOOGLE_API_KEY)


app = FastAPI(title="Document Assistant API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def get_file_extension(filename: str) -> str:
    # Robustly extract extension even on Windows paths like C:\path\file.jpg
    try:
        _, ext = os.path.splitext(filename)
        return (ext or "").lower()
    except Exception:
        return ""


def read_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    texts: List[str] = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(texts).strip()


def read_docx(file_bytes: bytes) -> str:
    file_like = io.BytesIO(file_bytes)
    doc = DocxDocument(file_like)
    return "\n".join(p.text for p in doc.paragraphs).strip()


def read_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text: List[str] = []
    for slide in prs.slides:
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_lines.append(getattr(shape, "text") or "")
        slides_text.append("\n".join(slide_lines))
    return "\n\n".join(slides_text).strip()


def read_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except Exception:
            continue
    return ""


def ocr_image_with_tesseract(img: Image.Image) -> str:
    if not TESS_AVAILABLE:
        return ""
    try:
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def extract_text_from_image(file_bytes: bytes, filename: str) -> str:
    try:
        img = Image.open(io.BytesIO(file_bytes))
    except Exception:
        return ""
    text = ocr_image_with_tesseract(img)
    return text.strip()


def extract_text_from_file(upload: UploadFile, file_bytes: bytes) -> str:
    ext = get_file_extension(upload.filename or "")
    if ext in (".pdf",):
        return read_pdf(file_bytes)
    if ext in (".docx",):
        return read_docx(file_bytes)
    if ext in (".pptx",):
        return read_pptx(file_bytes)
    if ext in (".txt",):
        return read_txt(file_bytes)
    if ext in (".png", ".jpg", ".jpeg"):
        return extract_text_from_image(file_bytes, upload.filename or "image")
    return ""


def chunk_text(text: str, max_chars: int = 20000) -> List[str]:
    chunks: List[str] = []
    if not text:
        return chunks
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        chunks.append(text[start:end])
        start = end
    return chunks


def ensure_gemini_configured() -> None:
    if not GOOGLE_API_KEY:
        raise RuntimeError("GOOGLE_API_KEY not set in environment")


def gemini_transcribe_audio(audio_bytes: bytes, mime_type: str) -> str:
    ensure_gemini_configured()
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = "Transcribe the audio faithfully. Return only the transcription."
    # Inline audio as input part
    audio_part = {"mime_type": mime_type, "data": audio_bytes}
    resp = model.generate_content([prompt, audio_part])
    return (resp.text or "").strip()


def gemini_ocr_image_if_needed(image_bytes: bytes, fallback_text: str) -> str:
    if fallback_text:
        return fallback_text
    try:
        ensure_gemini_configured()
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt = (
            "Extract all visible text from this image. "
            "Return plain text only."
        )
        image_part = {"mime_type": "image/png", "data": image_bytes}
        resp = model.generate_content([prompt, image_part])
        return (resp.text or "").strip()
    except Exception:
        return fallback_text


def build_multilingual_system_prompt(target_lang: str, has_document: bool = False) -> str:
    # Supported: English, Hindi, Malayalam, Marathi, Telugu, Bhojpuri
    if has_document:
        return (
            "You are a helpful, intelligent Document Assistant for students. "
            "Answer questions based on the provided document content when available. "
            "If the document doesn't contain relevant information, use your general knowledge to provide helpful answers. "
            "Always cite specific document snippets when available. "
            "Be accurate, concise, and educational. "
            "IMPORTANT: Format your response with proper structure:\n"
            "- Use **bold** for headings and important points\n"
            "- Use bullet points (•) for lists\n"
            "- Use numbered lists (1., 2., 3.) for steps\n"
            "- Use proper line breaks and spacing\n"
            "- Use indentation for sub-points\n"
            "- Use code blocks for technical content\n"
            "Use the user's requested language: " + target_lang + "."
        )
    else:
        return (
            "You are a comprehensive AI Knowledge Assistant for students. "
            "You can answer questions from any field including: "
            "General Knowledge, Geopolitics, Mathematics, Science, Entertainment, News, "
            "History, Geography, Literature, Technology, Sports, and more. "
            "Provide accurate, well-structured, and educational responses. "
            "Include relevant examples, facts, and explanations. "
            "IMPORTANT: Format your response with proper structure:\n"
            "- Use **bold** for headings and important points\n"
            "- Use bullet points (•) for lists\n"
            "- Use numbered lists (1., 2., 3.) for steps\n"
            "- Use proper line breaks and spacing\n"
            "- Use indentation for sub-points\n"
            "- Use code blocks for technical content\n"
            "- Use tables for structured data\n"
            "- Use > for quotes or important notes\n"
            "If you're unsure about recent events, mention that your knowledge has a cutoff date. "
            "Use the user's requested language: " + target_lang + "."
        )


def gemini_answer(document_text: str, user_query: str, language: str) -> str:
    ensure_gemini_configured()
    model = genai.GenerativeModel("gemini-1.5-flash")
    
    # Determine if we have document content
    has_document = bool(document_text and document_text.strip())
    system_prompt = build_multilingual_system_prompt(language, has_document)
    
    parts: List[object] = [system_prompt]
    
    # Add document content if available
    if has_document:
        doc_chunks = chunk_text(document_text, 20000)[:5]
        for idx, chunk in enumerate(doc_chunks):
            parts.append({"text": f"[Document Content - Chunk {idx+1}]\n{chunk}"})
    
    # Add user question
    parts.append({"text": f"[User Question]\n{user_query}"})
    
    resp = model.generate_content(parts)
    return (resp.text or "").strip()


@app.post("/ask")
async def ask(
    question: Optional[str] = Form(None),
    language: Optional[str] = Form("English"),
    file: Optional[UploadFile] = File(None),
    audio: Optional[UploadFile] = File(None),
):
    try:
        if not (file or audio or question):
            return JSONResponse(status_code=400, content={"error": "Provide a file or audio or question."})

        extracted_text = ""
        user_query = question or ""

        # If audio provided, transcribe to derive user query
        if audio is not None:
            audio_bytes = await audio.read()
            mime_type = audio.content_type or "audio/webm"
            if mime_type not in SUPPORTED_AUDIO_TYPES:
                # Try to infer from filename
                if audio.filename and audio.filename.lower().endswith(".wav"):
                    mime_type = "audio/wav"
            user_query = gemini_transcribe_audio(audio_bytes, mime_type)

        # If file provided, extract text (document or image)
        if file is not None:
            file_bytes = await file.read()
            ext = get_file_extension(file.filename or "")
            if ext in SUPPORTED_IMG_TYPES:
                extracted_text = extract_text_from_image(file_bytes, file.filename or "image")
                # Fallback to Gemini OCR if local OCR failed or Tesseract unavailable
                if not extracted_text:
                    extracted_text = gemini_ocr_image_if_needed(file_bytes, extracted_text)
            else:
                extracted_text = extract_text_from_file(file, file_bytes)

        if not user_query:
            return JSONResponse(status_code=400, content={"error": "Missing user question (text or audio)."})

        # Call Gemini to answer grounded on document text (if any)
        answer_text = gemini_answer(extracted_text, user_query, language or "English")

        return {
            "question": user_query,
            "answer": answer_text,
            "language": language or "English",
            "chars_in_document": len(extracted_text or ""),
        }

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.get("/")
async def root():
    return {"status": "ok", "service": "Document Assistant API"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="0.0.0.0", port=int(os.getenv("PORT", "8000")), reload=False)


